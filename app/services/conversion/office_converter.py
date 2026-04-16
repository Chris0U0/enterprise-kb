"""
Word / PPT → Markdown 转换器
使用 markitdown（微软官方开源库），标题/表格/列表格式保留最完整
"""
from __future__ import annotations

import logging
import tempfile
from pathlib import Path

logger = logging.getLogger(__name__)


async def convert_word(file_data: bytes, filename: str) -> tuple[str, int | None]:
    """
    Word (.docx/.doc) → Markdown
    Returns: (markdown_text, page_count)
    """
    try:
        from markitdown import MarkItDown

        with tempfile.NamedTemporaryFile(suffix=Path(filename).suffix, delete=False) as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name

        mid = MarkItDown()
        result = mid.convert(tmp_path)
        md_text = result.text_content

        # Word 没有天然的页码概念，page_count 设为 None
        page_count = _estimate_word_pages(md_text)

        return md_text, page_count

    except ImportError:
        logger.warning("markitdown 未安装，回退到 python-docx")
        return await _convert_word_fallback(file_data, filename)
    finally:
        try:
            Path(tmp_path).unlink(missing_ok=True)
        except Exception:
            pass


async def _convert_word_fallback(file_data: bytes, filename: str) -> tuple[str, int | None]:
    """回退方案：使用 python-docx"""
    import io
    from docx import Document as DocxDocument

    doc = DocxDocument(io.BytesIO(file_data))
    lines = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            lines.append("")
            continue

        style_name = para.style.name if para.style else ""

        # 标题映射
        if "Heading 1" in style_name:
            lines.append(f"# {text}")
        elif "Heading 2" in style_name:
            lines.append(f"## {text}")
        elif "Heading 3" in style_name:
            lines.append(f"### {text}")
        elif "Heading 4" in style_name:
            lines.append(f"#### {text}")
        elif "List" in style_name:
            lines.append(f"- {text}")
        else:
            lines.append(text)

    # 处理表格
    for table in doc.tables:
        lines.append("")
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
            if i == 0:
                lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
        lines.append("")

    md_text = "\n".join(lines)
    return md_text, _estimate_word_pages(md_text)


async def convert_ppt(file_data: bytes, filename: str) -> tuple[str, int | None]:
    """
    PPT (.pptx/.ppt) → Markdown
    每页转为一个 Markdown 章节
    """
    try:
        from markitdown import MarkItDown

        with tempfile.NamedTemporaryFile(suffix=Path(filename).suffix, delete=False) as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name

        mid = MarkItDown()
        result = mid.convert(tmp_path)
        md_text = result.text_content

        # 计算幻灯片数
        slide_count = md_text.count("# ") + md_text.count("## Slide")
        return md_text, max(slide_count, 1)

    except ImportError:
        logger.warning("markitdown 未安装，回退到 python-pptx")
        return await _convert_ppt_fallback(file_data, filename)
    finally:
        try:
            Path(tmp_path).unlink(missing_ok=True)
        except Exception:
            pass


async def _convert_ppt_fallback(file_data: bytes, filename: str) -> tuple[str, int | None]:
    """回退方案：使用 python-pptx"""
    import io
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_data))
    lines = []

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## 幻灯片 {i}")
        lines.append("")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # 标题形状
                if shape.shape_type and "TITLE" in str(shape.shape_type):
                    lines.append(f"### {shape.text.strip()}")
                else:
                    lines.append(shape.text.strip())
                lines.append("")

        lines.append("---")
        lines.append("")

    return "\n".join(lines), len(prs.slides)


def _estimate_word_pages(md_text: str) -> int | None:
    """估算 Word 文档页数 (约 300 汉字/页或 1500 字符/页)"""
    char_count = len(md_text)
    if char_count == 0:
        return None
    return max(1, char_count // 1500)
